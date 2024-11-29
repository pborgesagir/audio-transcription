import streamlit as st
from pptx import Presentation
from pptx.util import Inches
import fitz  # PyMuPDF
import speech_recognition as sr
from pydub import AudioSegment
import io
import ezdxf
import pandas as pd
import tempfile

# Set the page configuration
st.set_page_config(
    page_title='Conversão e Transcrição - SERCOM',
    layout='wide',
    page_icon="https://media.licdn.com/dms/image/C4D0BAQHXylmAyGyD3A/company-logo_200_200/0/1630570245289?e=2147483647&v=beta&t=Dxas2us5gteu0P_9mdkQBwJEyg2aoc215Vrk2phu7Bs",
    initial_sidebar_state='auto'
)

# Function to convert PDF to PowerPoint with each page as an image
def convert_pdf_to_ppt_with_fitz(pdf_file):
    pdf_document = fitz.open(stream=pdf_file.read(), filetype="pdf")
    presentation = Presentation()
    
    for page_number in range(len(pdf_document)):
        page = pdf_document.load_page(page_number)
        pix = page.get_pixmap()
        image_stream = io.BytesIO(pix.tobytes("png"))
        
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank layout
        slide.shapes.add_picture(image_stream, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))
    
    return presentation

# Function to transcribe audio
def transcribe_audio(audio_file):
    recognizer = sr.Recognizer()
    audio_segment = AudioSegment.from_file_using_temporary_files(audio_file, format='m4a')
    chunk_length_ms = 30000  # 30 seconds in milliseconds
    chunks = [audio_segment[i:i + chunk_length_ms] for i in range(0, len(audio_segment), chunk_length_ms)]
    full_text = ""
    
    for i, chunk in enumerate(chunks):
        with io.BytesIO() as audio_stream:
            chunk.export(audio_stream, format="wav")
            audio_stream.seek(0)
            with sr.AudioFile(audio_stream) as source:
                audio_data = recognizer.record(source)
                try:
                    text = recognizer.recognize_google(audio_data, language='pt-BR')
                    full_text += text + " "
                except sr.UnknownValueError:
                    st.write(f"Chunk {i+1} could not be transcribed.")
                except sr.RequestError as e:
                    st.write(f"Could not request results from Google Speech Recognition service; {e}")
    return full_text

# Function to extract locations and areas from DXF file
def extract_locations_from_dxf(dxf_file):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".dxf") as temp_file:
        temp_file.write(dxf_file.read())  # Write the uploaded file content to the temp file
        temp_file_path = temp_file.name
    
    # Read DXF file with ezdxf
    doc = ezdxf.readfile(temp_file_path)
    data = []
    for entity in doc.modelspace().query('TEXT MTEXT'):
        location_name = entity.plain_text() if hasattr(entity, 'plain_text') else entity.dxf.text
        area = entity.dxf.height if hasattr(entity.dxf, 'height') else None
        data.append({'Location Name': location_name, 'Area': area})
    
    return pd.DataFrame(data)

# Streamlit app layout
st.title('Conversão e Transcrição - SERCOM')
tab1, tab2, tab3 = st.tabs(["Conversão PDF para PPT", "Transcrição de Áudio", "DXF/DWG para CSV"])

# PDF to PPT Conversion
with tab1:
    st.subheader('Conversão de PDF para PPT')
    pdf_file = st.file_uploader("Faça o upload do arquivo em formato PDF", type=['pdf'])

    if pdf_file is not None:
        if st.button('Converter para PPT'):
            with st.spinner('Convertendo...'):
                try:
                    presentation = convert_pdf_to_ppt_with_fitz(pdf_file)
                    ppt_io = io.BytesIO()
                    presentation.save(ppt_io)
                    ppt_io.seek(0)
                    st.success("Conversão concluída 🥳")
                    st.download_button(
                        label="Baixar PPT",
                        data=ppt_io,
                        file_name="converted_presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                except Exception as e:
                    st.error(f"Error during conversion: {e}")

# Audio Transcription
with tab2:
    st.subheader('Transcrição de Áudio para Texto')
    audio_file = st.file_uploader("Faça o upload do arquivo em formato M4A", type=['m4a'])

    if audio_file is not None:
        st.audio(audio_file, format='audio/m4a')
        if st.button('Transcreva'):
            with st.spinner('Transcrevendo...'):
                try:
                    text = transcribe_audio(audio_file)
                    st.success("Transcrição concluída 🥳")
                    st.write(text)
                except Exception as e:
                    st.error(f"Error during transcription: {e}")

# DXF/DWG to CSV Conversion
with tab3:
    st.subheader('DXF/DWG para CSV')
    dwg_or_dxf_file = st.file_uploader("Faça o upload do arquivo em formato DXF ou DWG", type=['dxf', 'dwg'])

    if dwg_or_dxf_file is not None:
        file_extension = dwg_or_dxf_file.name.split('.')[-1].lower()
        if file_extension == 'dwg':
            st.error("Arquivos DWG não são suportados diretamente. Converta para DXF e tente novamente.")
        elif file_extension == 'dxf':
            if st.button('Extrair Dados e Baixar CSV'):
                with st.spinner('Processando...'):
                    try:
                        data = extract_locations_from_dxf(dwg_or_dxf_file)
                        csv_io = io.BytesIO()
                        data.to_csv(csv_io, index=False)
                        csv_io.seek(0)
                        st.success("Processamento concluído 🥳")
                        st.download_button(
                            label="Baixar CSV",
                            data=csv_io,
                            file_name="locations_and_areas.csv",
                            mime="text/csv"
                        )
                    except Exception as e:
                        st.error(f"Error during processing: {e}")
